"""Read source documents into plain text plus metadata.

Supported: .md, .txt, .pdf. Markdown/text files may start with a small header
block terminated by a '---' line:

    title: Gothic architecture
    source_url: https://...
    license: CC BY-SA 4.0
    ---

Those keys become chunk metadata; the header itself is not embedded.
"""

from __future__ import annotations

import hashlib
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

SUPPORTED_SUFFIXES = {".md", ".txt", ".pdf", ".docx", ".pptx"}
_HEADER_KEYS = {"title", "source_url", "license"}


@dataclass
class LoadedDocument:
    source: str  # path relative to the documents dir, used as the chunk key
    text: str
    sha256: str
    mtime: float
    metadata: dict[str, Any] = field(default_factory=dict)


def _parse_header(text: str) -> tuple[dict[str, Any], str]:
    lines = text.splitlines()
    metadata: dict[str, Any] = {}
    for index, line in enumerate(lines[:10]):
        stripped = line.strip()
        if stripped == "---":
            return metadata, "\n".join(lines[index + 1 :]).lstrip()
        if ":" in stripped:
            key, _, value = stripped.partition(":")
            key = key.strip().lower()
            if key in _HEADER_KEYS:
                metadata[key] = value.strip()
                continue
        break
    return {}, text


def _read_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(pages)


def _read_docx(path: Path) -> str:
    from docx import Document

    document = Document(str(path))
    blocks = [paragraph.text.strip() for paragraph in document.paragraphs]
    # Course handouts keep a lot of their content in tables.
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                blocks.append(" | ".join(cells))
    return "\n\n".join(block for block in blocks if block)


def _read_pptx(path: Path) -> str:
    from pptx import Presentation

    blocks: list[str] = []
    for number, slide in enumerate(Presentation(str(path)).slides, start=1):
        texts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if texts:
            # Slide headings carry the topic, so keep slides as their own blocks.
            blocks.append(f"[Slide {number}] " + "\n".join(texts))
    return "\n\n".join(blocks)


_READERS = {".pdf": _read_pdf, ".docx": _read_docx, ".pptx": _read_pptx}


class DocumentReadError(RuntimeError):
    """A file has a supported extension but cannot be parsed."""


def load_document(path: Path, documents_dir: Path) -> LoadedDocument | None:
    """Load one file. Returns None if it has no usable text."""
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_SUFFIXES:
        return None

    reader = _READERS.get(suffix)
    try:
        if reader:
            metadata, body = {}, reader(path)
        else:  # .md / .txt may carry a header block
            metadata, body = _parse_header(path.read_text(encoding="utf-8", errors="replace"))
    except Exception as exc:  # noqa: BLE001 - any parser failure is the same to us
        # Folders of real documents contain renamed, truncated and half-saved
        # files. One of them must not abort a whole ingestion run.
        raise DocumentReadError(f"{path.name}: {type(exc).__name__}: {exc}") from exc
    body = body.strip()
    if not body:
        return None

    metadata.setdefault("title", path.stem.replace("-", " ").replace("_", " ").title())
    stat = path.stat()
    return LoadedDocument(
        source=path.relative_to(documents_dir).as_posix(),
        text=body,
        sha256=hashlib.sha256(body.encode("utf-8")).hexdigest(),
        mtime=stat.st_mtime,
        metadata=metadata,
    )


def iter_documents(documents_dir: Path, on_error=None):
    """Yield every loadable document under documents_dir, sorted by path.

    Unreadable files are reported through on_error and skipped rather than
    ending the walk.
    """
    if not documents_dir.exists():
        return
    for path in sorted(documents_dir.rglob("*")):
        if not path.is_file() or path.suffix.lower() not in SUPPORTED_SUFFIXES:
            continue
        if path.name.startswith("~$"):  # Office lock files
            continue
        try:
            document = load_document(path, documents_dir)
        except DocumentReadError as exc:
            if on_error:
                on_error(exc)
            continue
        if document is not None:
            yield document
