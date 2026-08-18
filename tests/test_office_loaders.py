"""Loader coverage for the Office formats the course material arrives in."""

import pytest
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches

from src.loaders import load_document


def test_docx_paragraphs_and_tables_are_extracted(tmp_path):
    document = DocxDocument()
    document.add_paragraph("Reinforced concrete combines concrete and steel.")
    document.add_paragraph("")  # empty paragraphs should not produce blank blocks
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Material"
    table.cell(0, 1).text = "Compressive strength"
    table.cell(1, 0).text = "Concrete"
    table.cell(1, 1).text = "High"
    path = tmp_path / "materials.docx"
    document.save(str(path))

    loaded = load_document(path, tmp_path)

    assert "Reinforced concrete combines concrete and steel." in loaded.text
    assert "Material | Compressive strength" in loaded.text
    assert "Concrete | High" in loaded.text
    assert "\n\n\n" not in loaded.text
    assert loaded.metadata["title"] == "Materials"


def test_pptx_slides_are_extracted_with_numbers(tmp_path):
    presentation = Presentation()
    for index, heading in enumerate(["Green roofs", "Natural ventilation"], start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = heading
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
        box.text_frame.text = f"body text {index}"
    path = tmp_path / "climate.pptx"
    presentation.save(str(path))

    loaded = load_document(path, tmp_path)

    assert "[Slide 1]" in loaded.text and "[Slide 2]" in loaded.text
    assert "Green roofs" in loaded.text
    assert "Natural ventilation" in loaded.text
    assert loaded.text.index("Green roofs") < loaded.text.index("Natural ventilation")


def test_hebrew_text_survives_the_round_trip(tmp_path):
    document = DocxDocument()
    document.add_paragraph("אדריכלות מדברית מתאפיינת באקלים יבש ובקרינה גבוהה.")
    path = tmp_path / "desert.docx"
    document.save(str(path))

    loaded = load_document(path, tmp_path)
    assert "אדריכלות מדברית" in loaded.text


def test_empty_presentation_is_skipped(tmp_path):
    path = tmp_path / "blank.pptx"
    Presentation().save(str(path))
    assert load_document(path, tmp_path) is None


@pytest.mark.parametrize("name", ["notes.doc", "slides.ppt", "sheet.xlsx"])
def test_legacy_and_unsupported_office_formats_are_ignored(tmp_path, name):
    path = tmp_path / name
    path.write_bytes(b"binary")
    assert load_document(path, tmp_path) is None
